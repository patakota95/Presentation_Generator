import streamlit as st
import base64
from io import BytesIO
import tempfile
import os
import json
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

import google.generativeai as genai
from dotenv import load_dotenv

# Load environment variables
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
genai.configure(api_key=GEMINI_API_KEY)

# =============================================================================
# UCL Color Palette - Official colors only
# =============================================================================
UCL_COLORS = {
    'PURPLE': (80, 7, 120),        # #500778
    'GREEN': (82, 193, 82),        # #52C152
    'BLUE': (52, 198, 198),        # #34C6C6
    'YELLOW': (255, 202, 54),      # #FFCA36
    'PINK': (172, 20, 90),         # #AC145A
    'DARK_BLUE': (0, 34, 72),      # #002248
    'BLACK': (0, 0, 0),            # #000000
    'WHITE': (255, 255, 255)       # #FFFFFF
}

# Templates using official colors
TEMPLATES = {
    'ucl-blue': {
        'title_bg': UCL_COLORS['DARK_BLUE'],
        'title_text': UCL_COLORS['WHITE'],
        'content_bg': UCL_COLORS['WHITE'],
        'accent': UCL_COLORS['DARK_BLUE']
    },
    'ucl-purple': {
        'title_bg': UCL_COLORS['PURPLE'],
        'title_text': UCL_COLORS['WHITE'],
        'content_bg': UCL_COLORS['WHITE'],
        'accent': UCL_COLORS['PURPLE']
    },
    'ucl-green': {
        'title_bg': UCL_COLORS['GREEN'],
        'title_text': UCL_COLORS['BLACK'],
        'content_bg': UCL_COLORS['WHITE'],
        'accent': UCL_COLORS['GREEN']
    },
    'ucl-blue-vibrant': {
        'title_bg': UCL_COLORS['BLUE'],
        'title_text': UCL_COLORS['BLACK'],
        'content_bg': UCL_COLORS['WHITE'],
        'accent': UCL_COLORS['BLUE']
    },
    'ucl-pink': {
        'title_bg': UCL_COLORS['PINK'],
        'title_text': UCL_COLORS['WHITE'],
        'content_bg': UCL_COLORS['WHITE'],
        'accent': UCL_COLORS['PINK']
    }
}

# -----------------------------------------------------------------------------
# Helper: Convert hex color string to RGB tuple
# -----------------------------------------------------------------------------
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# -----------------------------------------------------------------------------
# Gemini 1.5 Flash text generation
# -----------------------------------------------------------------------------
def generate_content_with_gemini(bullet_points, audience, purpose):
    """
    Use a Gemini 1.5 Flash model to expand bullet points 
    into rich presentation content in JSON format.
    """
    try:
        bullet_points = bullet_points.strip()
        system_message = """You are an expert presentation content creator.
Your task is to transform simple bullet points into rich, detailed presentation content.
For each bullet point, provide:
1. An expanded main point that explains the concept clearly
2. 2-3 supporting sub-points with evidence, examples, or additional details
3. Brief speaker notes that guide the presenter on how to discuss this slide

Format your response as a JSON object with the following structure:
{
  "title": "A compelling title for the overall presentation",
  "points": [
    {
      "main_point": "Expanded version of bullet point 1",
      "sub_points": ["Supporting point 1", "Supporting point 2", "Supporting point 3"],
      "speaker_notes": "Notes for the presenter about this slide"
    }
  ],
  "conclusion": [
    "A key takeaway",
    "A practical application",
    "A call to action or next steps"
  ]
}

Make sure your content is professional, evidence-based, and tailored to the specified audience and purpose.
Do not add any explanation text outside the JSON structure.
"""
        user_message = f"""Transform these bullet points into detailed presentation content for {audience}, with the purpose to {purpose}:

{bullet_points}

Remember to format your response as JSON with a title, expanded points with sub-points, speaker notes, and conclusion points."""

        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content([system_message, user_message])
        return response.text

    except Exception as e:
        st.error(f"Error generating content with Gemini 1.5 Flash: {str(e)}")
        return None

# -----------------------------------------------------------------------------
# Parsing & Slide Generation
# -----------------------------------------------------------------------------
def parse_ai_response(response_text):
    """Parse the AI response into structured content for slides (JSON)."""
    try:
        json_match = re.search(r'({[\s\S]*})', response_text)
        if json_match:
            json_str = json_match.group(1)
            data = json.loads(json_str)
            return data
        else:
            return None
    except Exception as e:
        st.error(f"Error parsing AI response: {str(e)}")
        return None

def enhance_content_with_ai(bullet_points, audience, purpose, slide_count=5, ai_enabled=True):
    """Transform bullet points into structured slides with AI enhancement (via Gemini 1.5 Flash)."""
    # Extract bullet points from the text area
    points = [pt.strip().lstrip('- ') for pt in bullet_points.split('\n') 
              if pt.strip() and pt.strip().startswith('-')]
    if not points:
        return []
    
    ai_content = None
    if ai_enabled:
        # Generate content with Gemini 1.5 Flash
        ai_response = generate_content_with_gemini(bullet_points, audience, purpose)
        if ai_response:
            ai_content = parse_ai_response(ai_response)

    slides = []
    if ai_content and 'title' in ai_content:
        title = ai_content['title']
    else:
        main_words = " ".join(points[0].split()[:3]) if points else "UCL Presentation"
        title = f"UCL {main_words.title()}"
    subtitle = f"For {audience.title()}: {purpose.capitalize()}"
    slides.append({
        'title': title,
        'subtitle': subtitle,
        'content': [],
        'notes': f"Welcome. This presentation is designed for {audience} to {purpose}."
    })

    if ai_content and 'points' in ai_content and isinstance(ai_content['points'], list):
        for point_data in ai_content['points']:
            main_point = point_data.get('main_point', '')
            slide_title = main_point.split('.')[0] if main_point else "Key Point"
            content = [main_point]
            if 'sub_points' in point_data and point_data['sub_points']:
                content.extend(point_data['sub_points'])
            notes = point_data.get('speaker_notes', f"Discuss how this relates to {audience}.")
            slides.append({'title': slide_title, 'content': content, 'notes': notes})

        # Conclusion
        if 'conclusion' in ai_content and isinstance(ai_content['conclusion'], list):
            slides.append({
                'title': "Next Steps",
                'content': ai_content['conclusion'],
                'notes': "Summarize the presentation and open for questions."
            })
        else:
            slides.append({
                'title': "Next Steps",
                'content': ["Summary of key points", f"Application for {audience}", "Contact for further information"],
                'notes': "Summarize the presentation and open for questions."
            })
    else:
        # If AI is disabled or content is not valid, just create basic slides
        if len(points) <= slide_count - 2:
            for point in points:
                words = point.split()
                slide_title = " ".join(words[:min(3, len(words))]).title()
                slide_title = f"Focus on {slide_title}"
                slides.append({
                    'title': slide_title,
                    'content': [point],
                    'notes': f"Key point: {point}. Explain how this benefits {audience}."
                })
        else:
            points_per_slide = max(1, len(points) // (slide_count - 2))
            for i in range(0, len(points), points_per_slide):
                group = points[i:i+points_per_slide]
                words = group[0].split()
                slide_title = " ".join(words[:min(3, len(words))]).title()
                slides.append({
                    'title': slide_title,
                    'content': group,
                    'notes': f"Discuss these {len(group)} related points and their impact."
                })
        slides.append({
            'title': "Next Steps",
            'content': ["Summary of key points", f"Application for {audience}", "Contact for further information"],
            'notes': "Summarize the presentation and open for questions."
        })

    # Limit slides if needed
    if len(slides) > slide_count:
        slides = [slides[0]] + slides[1:slide_count-1] + [slides[-1]]
    return slides

def create_presentation(slides, template_name='ucl-blue', font_color_rgb=(0,0,0)):
    """
    Create a UCL-branded PowerPoint presentation, forcing a custom font color.
    - White slide background
    - Accent bars from the chosen template
    - Title slide vs. content slides
    """
    template = TEMPLATES.get(template_name, TEMPLATES['ucl-blue'])
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides):
        # Use a blank layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(0),
            top=Inches(0),
            width=prs.slide_width,
            height=Inches(0.5)
        )
        accent_bar.fill.solid()
        if i == 0:
            accent_bar.fill.fore_color.rgb = RGBColor(*template['title_bg'])
        else:
            accent_bar.fill.fore_color.rgb = RGBColor(*template['accent'])
        accent_bar.line.fill.background()

        if i == 0:
            # Title Slide
            title_box = slide.shapes.add_textbox(
                left=Inches(0.5), top=Inches(0.8),
                width=Inches(12), height=Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.clear()

            title_para = title_frame.paragraphs[0]
            title_para.text = slide_data['title']
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(*font_color_rgb)

            if slide_data.get('subtitle'):
                subtitle_para = title_frame.add_paragraph()
                subtitle_para.text = slide_data['subtitle']
                subtitle_para.font.size = Pt(28)
                subtitle_para.font.color.rgb = RGBColor(*font_color_rgb)

            # UCL Logo
            logo_box = slide.shapes.add_textbox(
                left=Inches(0.5), top=Inches(0.1),
                width=Inches(2), height=Inches(0.5)
            )
            logo_frame = logo_box.text_frame
            logo_frame.text = "UCL"
            logo_para = logo_frame.paragraphs[0]
            logo_para.font.bold = True
            logo_para.font.size = Pt(36)
            logo_para.font.color.rgb = RGBColor(*font_color_rgb)

            # Footer bar
            footer_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left=Inches(0),
                top=prs.slide_height - Inches(0.7),
                width=prs.slide_width,
                height=Inches(0.7)
            )
            footer_bar.fill.solid()
            footer_bar.fill.fore_color.rgb = RGBColor(*template['accent'])
            footer_bar.line.fill.background()

            # Speaker notes
            if 'notes' in slide_data and slide_data['notes'] and hasattr(slide, 'notes_slide'):
                slide.notes_slide.notes_text_frame.text = slide_data['notes']

        else:
            # Content Slide
            title_box = slide.shapes.add_textbox(
                left=Inches(0.5), top=Inches(0.7),
                width=Inches(12), height=Inches(1)
            )
            title_frame = title_box.text_frame
            title_frame.clear()
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_data['title'].upper()
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(*font_color_rgb)

            content_box = slide.shapes.add_textbox(
                left=Inches(0.5), top=Inches(1.7),
                width=Inches(12), height=Inches(5.0)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for point in slide_data['content']:
                p = content_frame.add_paragraph()
                p.text = point
                p.level = 0
                p.bullet = True
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(*font_color_rgb)

            if 'notes' in slide_data and slide_data['notes'] and hasattr(slide, 'notes_slide'):
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data['notes']

            # "UCL" branding
            branding_box = slide.shapes.add_textbox(
                left=Inches(0.5), top=Inches(0.1),
                width=Inches(1.5), height=Inches(0.5)
            )
            branding_frame = branding_box.text_frame
            branding_frame.text = "UCL"
            branding_para = branding_frame.paragraphs[0]
            branding_para.font.bold = True
            branding_para.font.size = Pt(28)
            branding_para.font.color.rgb = RGBColor(*font_color_rgb)

        # Slide number
        slide_number_box = slide.shapes.add_textbox(
            left=prs.slide_width - Inches(0.8),
            top=prs.slide_height - Inches(0.5),
            width=Inches(0.7),
            height=Inches(0.4)
        )
        sn_frame = slide_number_box.text_frame
        sn_frame.clear()
        sn_para = sn_frame.paragraphs[0]
        sn_para.text = str(i + 1)
        sn_para.alignment = PP_ALIGN.CENTER
        sn_para.font.size = Pt(14)
        sn_para.font.bold = True
        sn_para.font.color.rgb = RGBColor(*font_color_rgb)

    temp_dir = tempfile.gettempdir()
    file_path = os.path.join(temp_dir, f"UCL_Presentation_{template_name}.pptx")
    prs.save(file_path)
    return file_path

def create_preview_html(slides):
    """Generate a preview of slides as HTML for display."""
    html = ""
    for i, slide in enumerate(slides):
        if i == 0:
            html += f"""
            <div style="border: 1px solid #ccc; border-radius: 5px; padding: 15px; margin-bottom: 10px;">
                <h3 style="margin-top: 0;">Slide {i+1}: Title Slide</h3>
                <h2>{slide['title']}</h2>
                <h4 style="color: #666;">{slide.get('subtitle', '')}</h4>
                <p><em><strong>Speaker notes:</strong> {slide.get('notes', '')}</em></p>
            </div>
            """
        else:
            content_html = ""
            for point in slide['content']:
                content_html += f"<li>{point}</li>"
            html += f"""
            <div style="border: 1px solid #ccc; border-radius: 5px; padding: 15px; margin-bottom: 10px;">
                <h3 style="margin-top: 0;">Slide {i+1}: Content Slide</h3>
                <h2>{slide['title']}</h2>
                <ul style="margin-left: 20px;">
                    {content_html}
                </ul>
                <p><em><strong>Speaker notes:</strong> {slide.get('notes', '')}</em></p>
            </div>
            """
    return html

def generate_presentation(bullet_points, audience, purpose, template, slide_count, use_ai, font_color_hex):
    """
    Orchestrates the entire process:
    1. Enhances content with AI if enabled (Gemini 1.5 Flash)
    2. Builds the preview HTML
    3. Creates the PPTX file with forced white background and chosen font color
    """
    try:
        slides = enhance_content_with_ai(bullet_points, audience, purpose, slide_count, use_ai)
        if not slides:
            return "Error: Could not process bullet points. Make sure they start with '-'.", None

        preview = create_preview_html(slides)

        # Convert hex color to an RGB tuple
        font_color_rgb = hex_to_rgb(font_color_hex)

        # Create the PPTX file
        file_path = create_presentation(slides, template, font_color_rgb)
        return preview, file_path

    except Exception as e:
        error_message = f"Error generating presentation: {str(e)}"
        st.error(error_message)
        return error_message, None

# -----------------------------------------------------------------------------
# Streamlit UI
# -----------------------------------------------------------------------------
st.set_page_config(page_title="UCL Presentation Generator (Gemini 1.5 Flash)", layout="wide")

# Custom CSS for styling
st.markdown(
    """
    <style>
    .appview-container .main {
        background-color: #f8f9fa;
        padding: 20px;
    }
    h1, h2, h3 {
        color: #002248;
    }
    .stButton>button {
        background-color: #002248 !important;
        color: #ffffff !important;
        border-radius: 5px !important;
    }
    .stSlider>div>div>span>div[role="slider"] {
        background-color: #500778 !important;
    }
    .stSlider>div>div>div>div {
        background-color: #C9C5CB !important;
    }
    .streamlit-expanderHeader {
        font-size: 1rem;
        font-weight: 600;
        color: #500778;
    }
    </style>
    """,
    unsafe_allow_html=True
)

st.title("UCL Presentation Generator (Gemini 1.5 Flash)")
st.markdown("Turn simple bullet points into professional UCL-branded presentations with Gemini 1.5 Flash.")

# Instructions
with st.expander("How to Use"):
    st.write("""
1. Enter your bullet points in the text area (each starting with '-').
2. Select your target audience and the purpose of your presentation.
3. Choose a UCL-branded template style, set the number of slides, and select your desired font color.
4. Check **Use AI** to generate expanded content via Gemini 1.5 Flash.
5. Click **Generate Presentation** to preview and download your presentation.

**Gemini 1.5 Flash Setup**:
- Ensure you have a GEMINI_API_KEY set as environment variable.
- Ensure you have the `google-generativeai` and `python-dotenv` libraries installed.
- The prompt is a combination of a system message and your bullet points, returning a JSON structure.
""")

# Color palette reference
with st.expander("UCL Color Palette"):
    color_html = "<div style='display: flex; flex-wrap: wrap; gap: 10px;'>"
    for color_name, rgb in UCL_COLORS.items():
        hex_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        color_html += f"""
        <div style="width: 80px; text-align: center;">
            <div style="height: 40px; background-color: {hex_color}; margin-bottom: 5px; border-radius: 3px;"></div>
            <div style="font-size: 12px;">{color_name}</div>
        </div>
        """
    color_html += "</div>"
    st.markdown(color_html, unsafe_allow_html=True)

st.write("---")

# Example bullet points
example_bullets = """
- Places to visit in Geneva
- How much it will cost
"""

col_left, col_right = st.columns([1.2, 1.0], gap="large")

with col_left:
    st.subheader("Enter Bullet Points")
    st.info("Each bullet point should begin with a dash `-` on a new line.")
    bullet_points = st.text_area("Bullet Points", value=example_bullets, height=150)

with col_right:
    st.subheader("Presentation Details")
    audience = st.selectbox(
        "Select Audience",
        ["tourists", "students", "faculty", "researchers", "stakeholders", "business professionals", "general public"],
        index=0
    )
    purpose = st.text_input("Purpose of Presentation", value="provide travel information")
    template = st.selectbox("Select Template", list(TEMPLATES.keys()), index=0)
    slide_count = st.slider("Number of Slides", min_value=3, max_value=15, value=5, step=1)
    use_ai = st.checkbox("Use AI (Gemini 1.5 Flash Enhancement)", value=True)
    font_color_hex = st.color_picker("Select Font Color for Slide Text", "#000000")

st.write("---")

if st.button("Generate Presentation"):
    with st.spinner("Generating presentation via Gemini 1.5 Flash..."):
        preview, file_path = generate_presentation(
            bullet_points,
            audience,
            purpose,
            template,
            slide_count,
            use_ai,
            font_color_hex
        )
    st.markdown("## Presentation Preview")
    st.markdown(preview, unsafe_allow_html=True)

    if file_path and os.path.exists(file_path):
        with open(file_path, "rb") as f:
            file_bytes = f.read()
        st.download_button(
            label="Download Presentation",
            data=file_bytes,
            file_name="UCL_Presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )